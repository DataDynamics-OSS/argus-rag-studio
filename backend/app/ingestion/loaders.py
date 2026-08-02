# SPDX-License-Identifier: Apache-2.0
"""문서 로더 — 원본 바이트에서 평문 텍스트를 추출한다.

파일명 확장자로 포맷을 판별한다(``_LOADERS`` 레지스트리). 포맷을 추가할 때는
``if/elif`` 를 늘리는 대신 로더 함수를 작성해 ``_LOADERS`` 에 등록만 하면 된다.

바이너리 포맷(PDF/DOCX/XLSX/PPTX)은 선택 의존성을 사용하며 미설치 시 명확한 오류를
던진다. HTML/XML 은 stdlib ``html.parser`` 로 태그를 제거하고, 그 외 텍스트·설정·소스코드
계열은 디코딩만 한다. 레이아웃 인식·표 직렬화·OCR 은 ``parsers.py`` 의 전략에서 확장한다.
"""

import io
import logging
import zlib
from collections.abc import Callable
from html.parser import HTMLParser

logger = logging.getLogger(__name__)

# 평문 디코딩으로 충분한 문서/데이터 포맷
TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".log", ".rst"}

# 설정·마크업 계열 — 디코딩만으로 RAG 에 충분(2순위)
CONFIG_EXTENSIONS = {
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env", ".properties", ".tex",
}

# 소스코드 — 코드 RAG 용. 평문 디코딩.
CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".c", ".cc", ".cpp",
    ".h", ".hpp", ".cs", ".rb", ".php", ".sh", ".bash", ".sql", ".kt", ".swift",
    ".scala", ".r", ".pl",
}

# 디코딩만 하면 되는 모든 확장자(레지스트리 구성용)
PLAINTEXT_EXTENSIONS = TEXT_EXTENSIONS | CONFIG_EXTENSIONS | CODE_EXTENSIONS


class _HTMLTextExtractor(HTMLParser):
    """script/style 를 제외한 텍스트 노드만 모은다."""

    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip and data.strip():
            self._parts.append(data.strip())

    def text(self) -> str:
        return "\n".join(self._parts)


def _ext(filename: str) -> str:
    name = filename.lower()
    dot = name.rfind(".")
    return name[dot:] if dot >= 0 else ""


def _rows_to_markdown(rows: list[list[str]]) -> str:
    """행 리스트(문자열 셀)를 GitHub Markdown 표로 직렬화한다(첫 행을 헤더로)."""
    rows = [r for r in rows if any(cell for cell in r)]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header, body = rows[0], rows[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(lines)


def _load_text(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "cp949", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _load_html(data: bytes) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(_load_text(data))
    return parser.text()


def _load_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError("PDF 처리에는 pypdf 가 필요합니다. pip install pypdf") from e
    reader = PdfReader(io.BytesIO(data))
    pages = [(page.extract_text() or "") for page in reader.pages]
    return "\n\n".join(p.strip() for p in pages if p.strip())


def _load_docx(data: bytes) -> str:
    try:
        import docx  # python-docx
    except ImportError as e:
        raise RuntimeError("DOCX 처리에는 python-docx 가 필요합니다. pip install python-docx") from e
    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())


def _load_xlsx(data: bytes) -> str:
    """워크북의 각 시트를 ``# 시트명`` + Markdown 표로 직렬화한다(빈 셀/행 제거)."""
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise RuntimeError("XLSX 처리에는 openpyxl 이 필요합니다. pip install openpyxl") from e
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for ws in wb.worksheets:
            rows: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                cells = [("" if c is None else str(c)).strip().replace("\n", " ") for c in row]
                if any(cells):
                    rows.append(cells)
            table = _rows_to_markdown(rows)
            if table:
                parts.append(f"# {ws.title}\n\n{table}")
        return "\n\n".join(parts)
    finally:
        wb.close()


def _load_pptx(data: bytes) -> str:
    """슬라이드별 텍스트·표·발표자 노트를 ``## 슬라이드 N`` 단위로 추출한다."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("PPTX 처리에는 python-pptx 가 필요합니다. pip install python-pptx") from e
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                md = _rows_to_markdown(rows)
                if md:
                    texts.append(md)
        if slide.has_notes_slide:
            note = (slide.notes_slide.notes_text_frame.text or "").strip()
            if note:
                texts.append(f"(노트) {note}")
        if texts:
            parts.append(f"## 슬라이드 {idx}\n\n" + "\n".join(texts))
    return "\n\n".join(parts)


# ── HWPX (신형, OWPML = ZIP + XML) ────────────────────────────────────────────
def _xml_local(tag: str) -> str:
    """네임스페이스를 제거한 로컬 태그명(``{ns}p`` → ``p``)."""
    return tag.rsplit("}", 1)[-1]


def _hwpx_render(elem, notes: list[str], marks: dict | None = None) -> tuple[str, list[str]]:
    """``elem`` 의 (인라인 텍스트, 블록 리스트) 를 반환한다.

    표(``tbl``)는 Markdown 표 블록으로, 각주/미주(``footNote``/``endNote``)는 ``notes``
    에 라벨링해 적재한다(본문에 인라인 중복되지 않음). 그 외는 하위 ``t`` 를 이어붙인다.

    ``marks`` 가 주어지면(이미지 위치 마커 모드) ``binaryItemIDRef`` 를 가진 그림(``pic``)
    을 만날 때마다 등장 순번 카운터(``marks["n"]``)를 올리고, ``marks["layout"]`` 에 해당
    순번이 있으면 마커 블록(``marks["fmt"]``)을 그 자리에 끼운다 — 순번 축은 추출기의
    ``_hwpx_pic_media_map`` 과 동일(pre-order).
    """
    inline: list[str] = []
    blocks: list[str] = []
    for child in elem:
        name = _xml_local(child.tag)
        if name == "t":
            if child.text:
                inline.append(child.text)
        elif name == "tbl":
            md = _hwpx_table_md(child, notes, marks)
            if md:
                blocks.append(md)
        elif name in ("footNote", "endNote"):
            label = "각주" if name == "footNote" else "미주"
            ci, cb = _hwpx_render(child, notes, marks)
            txt = " ".join(s for s in [ci.strip(), *cb] if s).strip()
            if txt:
                notes.append(f"[{label}] {txt}")
        elif name == "p":
            ci, cb = _hwpx_render(child, notes, marks)
            if ci.strip():
                blocks.append(ci.strip())
            blocks.extend(cb)
        else:
            if marks is not None and name == "pic":
                ref = next(
                    (d.get("binaryItemIDRef") for d in child.iter() if d.get("binaryItemIDRef")),
                    None,
                )
                if ref:
                    marks["n"] += 1
                    for it in marks["layout"].get(marks["n"], ()):
                        marker = marks["fmt"].format(page=it["page"], index=it["index"])
                        # 표 셀 안 마커는 markdown 표를 깨뜨린다(주입 시 줄바꿈 삽입) —
                        # 표 렌더가 끝난 뒤 표 블록 뒤에 붙이도록 미룬다.
                        if marks.get("tbl_depth"):
                            marks.setdefault("deferred", []).append(marker)
                        else:
                            blocks.append(marker)
            ci, cb = _hwpx_render(child, notes, marks)
            inline.append(ci)
            blocks.extend(cb)
    return "".join(inline), blocks


def _hwpx_table_md(tbl, notes: list[str], marks: dict | None = None) -> str:
    """OWPML 표(``tbl`` → ``tr`` → ``tc``)를 Markdown 표로 직렬화한다.

    마커 모드에서는 표 안 그림의 마커를 미뤄뒀다가(최상위 표 렌더 종료 시) 표 블록
    바로 뒤에 줄 단위로 덧붙인다 — 셀 안에 넣으면 캡션 주입이 표를 깨뜨린다.
    """
    if marks is not None:
        marks["tbl_depth"] = marks.get("tbl_depth", 0) + 1
    rows: list[list[str]] = []
    for tr in tbl:
        if _xml_local(tr.tag) != "tr":
            continue
        cells: list[str] = []
        for tc in tr:
            if _xml_local(tc.tag) != "tc":
                continue
            ci, cb = _hwpx_render(tc, notes, marks)
            text = " ".join(s for s in [ci.strip(), *cb] if s).strip()
            cells.append(text.replace("\n", " "))
        if cells:
            rows.append(cells)
    md = _rows_to_markdown(rows)
    if marks is not None:
        marks["tbl_depth"] -= 1
        if marks["tbl_depth"] == 0 and marks.get("deferred"):
            md = "\n".join(([md] if md else []) + marks["deferred"])
            marks["deferred"] = []
    return md


def _hwpx_section_text(xml_bytes: bytes, marks: dict | None = None) -> str:
    """OWPML 섹션 XML 에서 본문·표·각주/미주를 추출한다(각주/미주는 말미에 첨부)."""
    import xml.etree.ElementTree as ET

    root = ET.fromstring(xml_bytes)
    notes: list[str] = []
    _, blocks = _hwpx_render(root, notes, marks)
    return "\n".join([b for b in blocks if b] + notes)


def load_hwpx_with_image_markers(data: bytes, layout: dict[int, list[dict]], fmt: str) -> str:
    """``_load_hwpx`` 와 동일한 본문 + 그림 등장 순번(``layout`` 키) 자리에 마커 삽입.

    ``fmt`` 는 마커 포맷 문자열(``[[IMG:p{page}:i{index}]]`` — image_classifier 정본을
    인자로 받아 역의존을 피한다). 카운터는 섹션 경계를 넘어 이어진다.
    """
    import zipfile

    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile as e:
        raise RuntimeError("올바른 HWPX 파일이 아닙니다(ZIP 컨테이너 손상).") from e
    marks = {"n": 0, "layout": layout, "fmt": fmt}
    with zf:
        names = sorted(
            n for n in zf.namelist()
            if n.startswith("Contents/section") and n.lower().endswith(".xml")
        )
        parts = [_hwpx_section_text(zf.read(n), marks) for n in names]
    return "\n".join(p for p in parts if p)


def _load_hwpx(data: bytes) -> str:
    """HWPX(ZIP) 컨테이너의 ``Contents/section*.xml`` 본문을 추출한다(외부 의존성 없음)."""
    import zipfile

    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile as e:
        raise RuntimeError("올바른 HWPX 파일이 아닙니다(ZIP 컨테이너 손상).") from e
    with zf:
        names = sorted(
            n for n in zf.namelist()
            if n.startswith("Contents/section") and n.lower().endswith(".xml")
        )
        parts = [_hwpx_section_text(zf.read(n)) for n in names]
    return "\n".join(p for p in parts if p)


# ── HWP 5.0 (구형 바이너리, OLE/CFB) ──────────────────────────────────────────
_HWPTAG_BEGIN = 0x10
_HWPTAG_PARA_TEXT = _HWPTAG_BEGIN + 51    # 67 — 문단 텍스트
_HWPTAG_CTRL_HEADER = _HWPTAG_BEGIN + 55  # 71 — 컨트롤 헤더(표/각주/미주 등)
_HWPTAG_LIST_HEADER = _HWPTAG_BEGIN + 56  # 72 — 리스트 헤더(표 셀·노트 본문 시작)
_HWPTAG_TABLE = _HWPTAG_BEGIN + 60        # 76 — 표 속성(행/열 수)

# 8 WCHAR(16바이트)을 차지하는 인라인·확장 컨트롤 문자(본문에서 건너뛴다)
_HWP_WIDE_CONTROLS = {1, 2, 3, 4, 5, 6, 7, 8, 9, 11, 12, 14, 15, 16, 17, 18, 19, 20, 21, 22, 23}


def _decode_hwp_para_text(rec: bytes) -> str:
    """PARA_TEXT 레코드(UTF-16LE + 인라인 컨트롤)를 평문으로 디코딩한다."""
    chars: list[str] = []
    i, m = 0, len(rec)
    while i + 2 <= m:
        code = int.from_bytes(rec[i:i + 2], "little")
        if code in _HWP_WIDE_CONTROLS:
            i += 16  # 인라인/확장 컨트롤 = 8 WCHAR
        elif code in (10, 13):  # 라인/문단 구분
            chars.append("\n")
            i += 2
        elif code < 32:  # 기타 char 컨트롤(1 WCHAR) — 무시
            i += 2
        else:
            chars.append(chr(code))
            i += 2
    return "".join(chars)


def _hwp_records(body: bytes):
    """해제된 섹션 스트림을 ``(tag, level, data)`` 레코드로 순회한다.

    레코드 헤더(UINT32 LE): tag_id(0~9) | level(10~19) | size(20~31).
    size==0xFFF 이면 다음 UINT32 가 실제 크기.
    """
    pos, n = 0, len(body)
    while pos + 4 <= n:
        header = int.from_bytes(body[pos:pos + 4], "little")
        pos += 4
        tag = header & 0x3FF
        level = (header >> 10) & 0x3FF
        size = (header >> 20) & 0xFFF
        if size == 0xFFF:
            if pos + 4 > n:
                break
            size = int.from_bytes(body[pos:pos + 4], "little")
            pos += 4
        yield tag, level, body[pos:pos + size]
        pos += size


def _hwp_ctrl_id(data: bytes) -> str:
    """CTRL_HEADER 의 ctrl_id(MAKE_4CHID, 리틀엔디안 → 역순)를 4글자 문자열로."""
    if len(data) < 4:
        return ""
    return data[:4][::-1].decode("latin-1", "ignore")


def _hwp_section_text(body: bytes) -> str:
    """섹션 레코드를 순회하며 본문·표·각주/미주를 추출한다(각주/미주는 말미에 첨부).

    컨트롤(CTRL_HEADER)의 ctrl_id 로 표(``tbl``)·각주(``fn``)·미주(``en``)를 식별하고,
    레코드의 ``level`` 로 컨트롤 하위 영역을 추적한다. 표는 셀 텍스트를 모아 열 수(nCols)
    기준으로 행을 묶어 Markdown 표로 직렬화한다(병합 셀은 근사). 본문 PARA_TEXT 는 줄로,
    각주/미주는 ``[각주]``/``[미주]`` 라벨로 적재한다. 그 외 컨트롤 내부 텍스트는 본문 취급.
    """
    lines: list[str] = []
    notes: list[str] = []
    stack: list[dict] = []  # 활성 컨트롤 컨텍스트(표/노트)

    def _emit_block(block: str) -> None:
        if not block:
            return
        if stack and stack[-1]["kind"] == "table":
            top = stack[-1]
            if top["cur"] is None:
                top["cur"] = []
            top["cur"].append(block)
        elif stack and stack[-1]["kind"] == "note":
            stack[-1]["text"].append(block)
        else:
            lines.append(block)

    def _finalize(ctx: dict) -> None:
        if ctx["kind"] == "table":
            if ctx.get("cur"):
                ctx["cells"].append(" ".join(ctx["cur"]).strip())
            ncols = ctx["ncols"] or max(len(ctx["cells"]), 1)
            rows = [
                [c.replace("\n", " ") for c in ctx["cells"][i:i + ncols]]
                for i in range(0, len(ctx["cells"]), ncols)
            ]
            _emit_block(_rows_to_markdown(rows))
        elif ctx["kind"] == "note":
            txt = " ".join(ctx["text"]).strip().replace("\n", " ")
            if txt:
                notes.append(f"[{ctx['label']}] {txt}")

    for tag, level, data in _hwp_records(body):
        while stack and level <= stack[-1]["level"]:
            _finalize(stack.pop())

        if tag == _HWPTAG_CTRL_HEADER:
            cid = _hwp_ctrl_id(data).strip()
            if cid == "tbl":
                stack.append({"kind": "table", "level": level, "ncols": 0, "cells": [], "cur": None})
            elif cid == "fn":
                stack.append({"kind": "note", "level": level, "label": "각주", "text": []})
            elif cid == "en":
                stack.append({"kind": "note", "level": level, "label": "미주", "text": []})
        elif tag == _HWPTAG_TABLE:
            if stack and stack[-1]["kind"] == "table" and len(data) >= 8:
                # UINT32 속성(4) + UINT16 nRows(2) + UINT16 nCols(2)
                stack[-1]["ncols"] = int.from_bytes(data[6:8], "little")
        elif tag == _HWPTAG_LIST_HEADER:
            if stack and stack[-1]["kind"] == "table":
                top = stack[-1]
                if top["cur"] is not None:
                    top["cells"].append(" ".join(top["cur"]).strip())
                top["cur"] = []
        elif tag == _HWPTAG_PARA_TEXT:
            text = _decode_hwp_para_text(data)
            if stack and stack[-1]["kind"] == "table":
                top = stack[-1]
                if top["cur"] is None:
                    top["cur"] = []
                top["cur"].append(text)
            elif stack and stack[-1]["kind"] == "note":
                stack[-1]["text"].append(text)
            else:
                lines.append(text)

    while stack:
        _finalize(stack.pop())

    body_text = "\n".join(s.strip() for s in lines if s.strip())
    if notes:
        body_text = (body_text + "\n" if body_text else "") + "\n".join(notes)
    return body_text


def _load_hwp(data: bytes) -> str:
    """HWP 5.0(OLE) 의 ``BodyText/Section*`` 을 해제·파싱한다(olefile 필요). 암호화 시 오류."""
    try:
        import olefile
    except ImportError as e:
        raise RuntimeError("HWP 처리에는 olefile 이 필요합니다. pip install olefile") from e

    buf = io.BytesIO(data)
    if not olefile.isOleFile(buf):
        raise RuntimeError("올바른 HWP 5.0 파일이 아닙니다(구형 3.0 또는 손상 가능).")
    ole = olefile.OleFileIO(buf)
    try:
        header = ole.openstream("FileHeader").read()
        if header[:17] != b"HWP Document File":
            raise RuntimeError("HWP 시그니처가 없습니다.")
        flags = int.from_bytes(header[36:40], "little")
        compressed = bool(flags & 0x01)
        if flags & 0x02:
            raise RuntimeError("암호화/DRM HWP 는 지원하지 않습니다(비밀번호 필요).")

        def _section_no(entry: list[str]) -> int:
            digits = "".join(c for c in entry[1] if c.isdigit())
            return int(digits) if digits else 0

        sections = sorted(
            (e for e in ole.listdir()
             if len(e) == 2 and e[0] == "BodyText" and e[1].lower().startswith("section")),
            key=_section_no,
        )
        parts: list[str] = []
        for entry in sections:
            raw = ole.openstream(entry).read()
            body = zlib.decompress(raw, -15) if compressed else raw
            parts.append(_hwp_section_text(body))
    finally:
        ole.close()
    return "\n".join(p for p in parts if p)


# 확장자 → 로더 레지스트리. 포맷 추가 시 여기 한 줄만 등록하면 된다.
_LOADERS: dict[str, Callable[[bytes], str]] = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".xlsx": _load_xlsx,
    ".pptx": _load_pptx,
    ".hwp": _load_hwp,
    ".hwpx": _load_hwpx,
    ".html": _load_html,
    ".htm": _load_html,
    ".xml": _load_html,
    "": _load_text,  # 확장자 없는 파일은 평문으로 처리
    **{ext: _load_text for ext in PLAINTEXT_EXTENSIONS},
}

# 지원 확장자 목록(UI accept·문서·테스트에서 단일 소스로 참조)
SUPPORTED_EXTENSIONS = sorted(e for e in _LOADERS if e)


def extract_text(filename: str, data: bytes) -> str:
    """파일명/바이트에서 평문 텍스트를 추출한다.

    Raises:
        RuntimeError: 미지원 포맷이거나 선택 의존성이 없을 때.
    """
    ext = _ext(filename)
    loader = _LOADERS.get(ext)
    if loader is None:
        raise RuntimeError(f"지원하지 않는 파일 형식입니다: {ext or '(확장자 없음)'}")

    text = loader(data).strip()
    if not text:
        raise RuntimeError("추출된 텍스트가 비어 있습니다(스캔본 PDF 등 OCR 필요 가능).")
    return text
