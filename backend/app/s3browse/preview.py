# SPDX-License-Identifier: Apache-2.0
"""오피스/데이터 파일 서버사이드 미리보기 — bytes → 구조화 데이터.

parquet/xlsx/xls → 테이블, docx/pptx → 문서(HTML). 프런트 뷰어(XlsxViewer/ParquetViewer/
DocxViewer/PptxViewer)가 기대하는 형태를 그대로 반환한다. 파서는 동기라 호출부에서 to_thread 로 오프로드.
"""

from __future__ import annotations

import html
import io
from datetime import datetime

from app.core.config import settings

# 미리보기 최대 행수 폴백 — 실제 값은 settings.preview_max_rows.
MAX_PREVIEW_ROWS = 1000


def _max_rows(v: int | None) -> int:
    return v if v is not None else settings.preview_max_rows


def _serialize_value(val):
    """셀 값을 JSON 직렬화 가능한 타입으로 — datetime/Decimal/bytes 등."""
    if val is None or isinstance(val, (int, float, str, bool)):
        return val
    if isinstance(val, datetime):
        return val.isoformat()
    if isinstance(val, bytes):
        return val.hex()
    return str(val)


def preview_parquet(data: bytes, max_rows: int | None = None) -> dict:
    import pyarrow.parquet as pq

    max_rows = _max_rows(max_rows)
    pf = pq.ParquetFile(io.BytesIO(data))
    total_rows = pf.metadata.num_rows
    columns = list(pf.schema_arrow.names)

    table = pf.read_row_groups(list(range(pf.metadata.num_row_groups)))
    if table.num_rows > max_rows:
        table = table.slice(0, max_rows)

    rows: list[list] = []
    for batch in table.to_batches():
        cols = [batch.column(i).to_pylist() for i in range(batch.num_columns)]
        for ri in range(batch.num_rows):
            rows.append([_serialize_value(cols[ci][ri]) for ci in range(len(cols))])

    return {
        "format": "parquet", "columns": columns, "rows": rows,
        "total_rows": total_rows, "sheet_names": [], "active_sheet": "",
    }


def preview_xlsx(data: bytes, sheet: str | None = None, max_rows: int | None = None, ext: str = "xlsx") -> dict:
    import openpyxl

    max_rows = _max_rows(max_rows)
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheet_names = wb.sheetnames
    active = sheet if sheet and sheet in sheet_names else sheet_names[0]
    ws = wb[active]

    rows: list[list] = []
    columns: list[str] = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i == 0:
            columns = [str(c) if c is not None else f"col_{j}" for j, c in enumerate(row)]
            continue
        if i > max_rows:
            break
        rows.append([_serialize_value(c) for c in row])

    total_rows = ws.max_row - 1 if ws.max_row else 0
    wb.close()
    return {
        "format": ext or "xlsx", "columns": columns, "rows": rows,
        "total_rows": total_rows, "sheet_names": sheet_names, "active_sheet": active,
    }


def preview_docx(data: bytes) -> dict:
    import mammoth

    result = mammoth.convert_to_html(io.BytesIO(data))
    return {"format": "docx", "html": result.value, "slides": None}


def preview_pptx(data: bytes) -> dict:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(cell.text.strip() for cell in row.cells))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({"slide_number": i, "texts": texts, "notes": notes})

    # 슬라이드 텍스트는 사용자 콘텐츠라 escape 후 HTML 구성(XSS 방지)
    parts: list[str] = []
    for s in slides:
        parts.append(f'<div class="slide"><h3>Slide {s["slide_number"]}</h3>')
        for t in s["texts"]:
            parts.append(f"<p>{html.escape(t)}</p>")
        if s["notes"]:
            parts.append(f'<blockquote class="notes">{html.escape(s["notes"])}</blockquote>')
        parts.append("</div><hr/>")

    return {"format": "pptx", "html": "\n".join(parts), "slides": slides}
