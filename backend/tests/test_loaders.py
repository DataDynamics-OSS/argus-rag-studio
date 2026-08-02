# SPDX-License-Identifier: Apache-2.0
"""문서 로더(loaders.py) 테스트 — 확장자 레지스트리 라우팅 + 포맷별 추출.

바이너리 포맷(xlsx/pptx)은 픽스처를 인메모리로 생성하므로 해당 라이브러리가 있을 때만
실행한다(없으면 skip). 텍스트/설정/소스코드 계열과 미지원 포맷 처리는 의존성 없이 검증한다.
"""

import importlib.util
import io
import zipfile

import pytest

from app.ingestion.loaders import (
    PLAINTEXT_EXTENSIONS,
    SUPPORTED_EXTENSIONS,
    _hwp_section_text,
    extract_text,
)


def _has(mod: str) -> bool:
    return importlib.util.find_spec(mod) is not None


# ── 레지스트리 라우팅 ──────────────────────────────────────────────────────────
def test_unsupported_format_raises():
    with pytest.raises(RuntimeError, match="지원하지 않는"):
        extract_text("archive.zip", b"PK\x03\x04")


def test_empty_extraction_raises():
    with pytest.raises(RuntimeError, match="비어 있습니다"):
        extract_text("blank.txt", b"   \n  ")


def test_no_extension_treated_as_text():
    assert "hello" in extract_text("README", b"hello world")


def test_supported_extensions_listing():
    # accept/문서가 참조하는 단일 소스 — 신규 포맷이 빠짐없이 포함되는지
    assert {".pdf", ".docx", ".xlsx", ".pptx", ".xml", ".hwp", ".hwpx"} <= set(SUPPORTED_EXTENSIONS)


# ── 텍스트/설정/소스코드 계열(의존성 없음) ────────────────────────────────────
@pytest.mark.parametrize("ext", [".yaml", ".toml", ".ini", ".tex", ".py", ".sql", ".go"])
def test_plaintext_family_extensions(ext):
    assert ext in PLAINTEXT_EXTENSIONS
    out = extract_text(f"file{ext}", "key: 값 RAG".encode())
    assert "RAG" in out


def test_xml_strips_tags():
    out = extract_text("doc.xml", b"<root><title>RAG</title></root>")
    assert "RAG" in out and "<title>" not in out


# ── xlsx (openpyxl) ────────────────────────────────────────────────────────────
@pytest.mark.skipif(not _has("openpyxl"), reason="openpyxl 미설치")
def test_xlsx_sheet_to_markdown():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Quarter", "Revenue"])
    ws.append(["Q1", 100])
    ws.append(["Q2", 150])
    buf = io.BytesIO()
    wb.save(buf)

    out = extract_text("report.xlsx", buf.getvalue())
    assert "# Sales" in out
    assert "| Quarter | Revenue |" in out
    assert "| Q1 | 100 |" in out


# ── pptx (python-pptx) ─────────────────────────────────────────────────────────
@pytest.mark.skipif(not _has("pptx"), reason="python-pptx 미설치")
def test_pptx_slide_text():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 제목만
    slide.shapes.title.text = "RAG 스튜디오 소개"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "임베딩 전 파싱"
    buf = io.BytesIO()
    prs.save(buf)

    out = extract_text("deck.pptx", buf.getvalue())
    assert "## 슬라이드 1" in out
    assert "RAG 스튜디오 소개" in out
    assert "임베딩 전 파싱" in out


# ── hwpx (ZIP + OWPML, 의존성 없음) ───────────────────────────────────────────
_HWPX_OPEN = (
    '<?xml version="1.0" encoding="UTF-8"?>'
    '<hs:sec xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section"'
    ' xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">'
)


def _cell(text: str) -> str:
    return f"<hp:tc><hp:subList><hp:p><hp:run><hp:t>{text}</hp:t></hp:run></hp:p></hp:subList></hp:tc>"


def _make_hwpx(body: str) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("Contents/section0.xml", _HWPX_OPEN + body + "</hs:sec>")
    return buf.getvalue()


def test_hwpx_extracts_paragraphs():
    out = extract_text("doc.hwpx", _make_hwpx(
        "<hp:p><hp:run><hp:t>안녕하세요 RAG 스튜디오</hp:t></hp:run></hp:p>"
        "<hp:p><hp:run><hp:t>두 번째 문단</hp:t></hp:run></hp:p>"
    ))
    assert "안녕하세요 RAG 스튜디오" in out
    assert "두 번째 문단" in out
    assert "<hp:t>" not in out  # 태그 노출 없음


def test_hwpx_table_to_markdown():
    body = (
        "<hp:p><hp:run><hp:t>본문 시작</hp:t></hp:run></hp:p>"
        "<hp:p><hp:run><hp:tbl>"
        f"<hp:tr>{_cell('제목')}{_cell('값')}</hp:tr>"
        f"<hp:tr>{_cell('매출')}{_cell('100')}</hp:tr>"
        "</hp:tbl></hp:run></hp:p>"
    )
    out = extract_text("doc.hwpx", _make_hwpx(body))
    assert "본문 시작" in out
    assert "| 제목 | 값 |" in out
    assert "| 매출 | 100 |" in out


def test_hwpx_footnote_labeled_at_end():
    body = (
        "<hp:p><hp:run><hp:t>각주 참조</hp:t>"
        "<hp:footNote><hp:subList><hp:p><hp:run><hp:t>각주 본문</hp:t></hp:run></hp:p></hp:subList></hp:footNote>"
        "</hp:run></hp:p>"
    )
    out = extract_text("doc.hwpx", _make_hwpx(body))
    assert "각주 참조" in out
    assert "[각주] 각주 본문" in out
    # 본문 줄이 각주보다 앞에 와야 함
    assert out.index("각주 참조") < out.index("[각주]")


def test_hwpx_bad_zip_raises():
    with pytest.raises(RuntimeError, match="HWPX"):
        extract_text("doc.hwpx", b"not a zip")


# ── hwp (5.0 바이너리 레코드 파서 — OLE 없이 핵심 로직 단위 검증) ─────────────
def _para_text_record(text: str, *, wide_control: bool = False, para_break: bool = False) -> bytes:
    """PARA_TEXT(tag 67) 레코드 1개를 합성한다(UTF-16LE + 선택적 인라인 컨트롤)."""
    payload = text.encode("utf-16-le")
    if wide_control:  # 인라인 컨트롤(코드 4) = 8 WCHAR = 16바이트 → 본문에서 스킵돼야 함
        payload += (4).to_bytes(2, "little") + b"\x00" * 14
    if para_break:  # 코드 13 = 문단 구분 → 개행
        payload += (13).to_bytes(2, "little")
    tag, level, size = 67, 0, len(payload)
    header = (tag & 0x3FF) | ((level & 0x3FF) << 10) | ((size & 0xFFF) << 20)
    return header.to_bytes(4, "little") + payload


def test_hwp_section_parser_decodes_text():
    section = _para_text_record("가A나")
    assert _hwp_section_text(section) == "가A나"


def test_hwp_section_parser_skips_inline_controls():
    # "가" + 인라인 컨트롤(16바이트) + "나" → 컨트롤은 사라지고 텍스트만 남음
    section = _para_text_record("가나", wide_control=True)
    assert _hwp_section_text(section) == "가나"


def test_hwp_section_parser_handles_para_break():
    section = _para_text_record("첫줄", para_break=True) + _para_text_record("둘째줄")
    assert _hwp_section_text(section) == "첫줄\n둘째줄"


# ── hwp 표·각주 레코드 합성 ───────────────────────────────────────────────────
def _rec(tag: int, level: int, payload: bytes) -> bytes:
    size = len(payload)
    header = (tag & 0x3FF) | ((level & 0x3FF) << 10) | ((size & 0xFFF) << 20)
    return header.to_bytes(4, "little") + payload


def _para(text: str, level: int) -> bytes:
    return _rec(67, level, text.encode("utf-16-le"))


def _ctrl(level: int, cid: str) -> bytes:
    # MAKE_4CHID: 문자열을 역순으로 저장(리틀엔디안)
    return _rec(71, level, bytes(reversed(cid.encode("latin-1"))))


def _table(level: int, nrows: int, ncols: int) -> bytes:
    payload = (0).to_bytes(4, "little") + nrows.to_bytes(2, "little") + ncols.to_bytes(2, "little")
    return _rec(76, level, payload)


def _list_header(level: int) -> bytes:
    return _rec(72, level, b"\x00" * 6)


def test_hwp_table_reconstructed_to_markdown():
    # 컨트롤(L) → 표(L+1) → 셀 4개(2x2): 제목/값/매출/100 → 본문(L) 으로 표 종료
    section = (
        _ctrl(1, "tbl ")
        + _table(2, 2, 2)
        + _list_header(2) + _para("제목", 4)
        + _list_header(2) + _para("값", 4)
        + _list_header(2) + _para("매출", 4)
        + _list_header(2) + _para("100", 4)
        + _para("본문 끝", 1)
    )
    out = _hwp_section_text(section)
    assert "| 제목 | 값 |" in out
    assert "| 매출 | 100 |" in out
    assert "본문 끝" in out


def test_hwp_footnote_labeled_at_end():
    section = (
        _para("본문1", 1)
        + _ctrl(1, "fn  ")
        + _para("각주 내용", 3)
        + _para("본문2", 1)
    )
    out = _hwp_section_text(section)
    assert out.startswith("본문1\n본문2")
    assert "[각주] 각주 내용" in out
    assert out.index("본문2") < out.index("[각주]")
