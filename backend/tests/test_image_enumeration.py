# SPDX-License-Identifier: Apache-2.0
"""인제스천 이미지 열거(포맷별 추출 위임) 테스트 — design/image-content-indexing.md Step 1.

이미지 열거가 imaging.extractors 로 위임되어 PDF 외 포맷(docx/pptx 등)도 분류 대상이
되는지, 상한 초과 시 큰 이미지 우선 선별이 순서를 유지하는지 검증한다. VLM 호출은
여기서 다루지 않는다(열거는 순수 로직).
"""

import io

from PIL import Image

from app.ingestion.image_classifier import _cap_largest, _enumerate_images_for_classify


def _png_bytes(w: int, h: int, color=(200, 30, 30)) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", (w, h), color).save(buf, format="PNG")
    return buf.getvalue()


def _docx_with_image() -> bytes:
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_paragraph("이미지 포함 문서")
    doc.add_picture(io.BytesIO(_png_bytes(320, 200)), width=Inches(2))
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _pptx_with_image() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(io.BytesIO(_png_bytes(400, 240)), Inches(1), Inches(1))
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ── 포맷별 열거(기존 PDF 전용 → 전 포맷) ─────────────────────────────────────

def test_docx_images_enumerated():
    items = _enumerate_images_for_classify("보고서.docx", _docx_with_image())
    assert len(items) == 1
    it = items[0]
    assert it["png"] and it["width"] == 320 and it["height"] == 200
    # Step 3: 문단 앵커 매핑 — page=문단 번호(텍스트 문단 1 + 그림 문단 2), unit="para".
    assert it["page"] == 2 and it["unit"] == "para"
    assert it["locator"].startswith("para 2 · word/media/")


def test_pptx_images_enumerated():
    items = _enumerate_images_for_classify("발표.pptx", _pptx_with_image())
    assert len(items) == 1
    # Step 3: 슬라이드 관계 매핑 — page=슬라이드 번호, unit="slide", locator 에 슬라이드 표기.
    it = items[0]
    assert it["page"] == 1 and it["unit"] == "slide"
    assert it["locator"].startswith("slide 1 · ppt/media/")


def test_single_image_file():
    items = _enumerate_images_for_classify("사진.png", _png_bytes(300, 300))
    assert len(items) == 1 and items[0]["png"]


def test_min_pixels_filter_applies():
    # 최소 크기(기본 64px) 미만 아이콘은 제외 — extractors 공통 필터.
    items = _enumerate_images_for_classify("아이콘.png", _png_bytes(16, 16))
    assert items == []


def test_docx_thumbnail_not_treated_as_image():
    # python-docx 기본 템플릿의 docProps/thumbnail.jpeg(미리보기 스냅샷)는 본문 이미지가 아니다 —
    # 이미지 없는 문서에 유령 이미지가 분류/주입되던 회귀 방지.
    from docx import Document
    out = io.BytesIO()
    Document().save(out)  # 그림 없음(썸네일만 존재)
    assert _enumerate_images_for_classify("빈문서.docx", out.getvalue()) == []


def test_unsupported_format_empty():
    assert _enumerate_images_for_classify("data.csv", b"a,b\n1,2") == []


# ── PPTX 슬라이드 매핑(Step 3) ───────────────────────────────────────────────

def _pptx_two_slides() -> bytes:
    """슬라이드 1·3에만 이미지가 있는 3장짜리 pptx(슬라이드 2 는 텍스트만)."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.add_picture(io.BytesIO(_png_bytes(400, 240)), Inches(1), Inches(1))
    prs.slides.add_slide(prs.slide_layouts[5])  # 이미지 없음
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.add_picture(io.BytesIO(_png_bytes(500, 300, (30, 30, 200))), Inches(1), Inches(1))
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def test_pptx_slide_numbers_mapped_per_slide():
    items = _enumerate_images_for_classify("발표.pptx", _pptx_two_slides())
    assert len(items) == 2
    assert sorted(it["page"] for it in items) == [1, 3]
    assert all(it["unit"] == "slide" for it in items)


def test_pptx_image_layout_groups_by_slide():
    from app.ingestion.image_classifier import media_image_layout

    layout = media_image_layout("발표.pptx", _pptx_two_slides())
    assert set(layout.keys()) == {1, 3}
    # index 는 열거(index)와 동일 규약 — (page, index) 가 items 와 1:1.
    for slide, its in layout.items():
        assert all(it["page"] == slide for it in its)


def test_append_slide_markers_in_section_and_tail():
    from app.ingestion.image_classifier import append_slide_markers

    text = "## 슬라이드 1\n\n첫 장 본문\n\n## 슬라이드 2\n\n둘째 장 본문"
    layout = {
        1: [{"index": 0, "page": 1}],
        3: [{"index": 1, "page": 3}],  # 본문 섹션 없는 슬라이드(이미지 전용)
    }
    out = append_slide_markers(text, layout)
    # 슬라이드 1 마커는 해당 섹션 안(다음 헤더 이전)에.
    assert out.index("첫 장 본문") < out.index("[[IMG:p1:i0]]") < out.index("## 슬라이드 2")
    # 섹션 없는 슬라이드 3 마커는 문서 말미에.
    assert out.rstrip().endswith("[[IMG:p3:i1]]")
    # 레이아웃이 비면 원문 그대로.
    assert append_slide_markers(text, {}) == text


# ── DOCX 문단 매핑(Step 3) ───────────────────────────────────────────────────

def _docx_two_images() -> bytes:
    """문단 1(텍스트)·2(그림)·3(텍스트)·4(그림) 구조의 docx."""
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_paragraph("서론 문단")
    doc.add_picture(io.BytesIO(_png_bytes(320, 200)), width=Inches(2))
    doc.add_paragraph("본론 문단")
    doc.add_picture(io.BytesIO(_png_bytes(400, 240, (30, 30, 200))), width=Inches(2))
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def test_docx_para_numbers_mapped():
    items = _enumerate_images_for_classify("보고서.docx", _docx_two_images())
    assert len(items) == 2
    assert sorted(it["page"] for it in items) == [2, 4]
    assert all(it["unit"] == "para" for it in items)


def test_docx_parse_inserts_markers_at_para():
    from app.ingestion.parsers import parse_document

    text = parse_document("보고서.docx", _docx_two_images(), "text", insert_image_markers=True)
    lines = text.split("\n")
    # 문단 축 그대로: 텍스트 문단 사이의 그림 문단 위치에 마커가 놓인다.
    i_intro, i_body = lines.index("서론 문단"), lines.index("본론 문단")
    assert any("[[IMG:p2:" in ln for ln in lines[i_intro + 1:i_body])
    assert any("[[IMG:p4:" in ln for ln in lines[i_body + 1:])
    # 마커 없이 파싱하면(미리보기 등) 로더 원문과 동일.
    plain = parse_document("보고서.docx", _docx_two_images(), "text")
    assert "[[IMG" not in plain and "서론 문단" in plain


# ── HTML img 순번 매핑(Step 3) ───────────────────────────────────────────────

def _html_with_images() -> bytes:
    """data URI img 3개(2번째는 min_pixels 미만) — 등장 순번 축 검증용."""
    import base64

    def uri(w, h):
        return "data:image/png;base64," + base64.b64encode(_png_bytes(w, h)).decode()

    return (
        f'<html><body><p>첫 문단</p><img src="{uri(320, 200)}">'
        f'<p>둘째 문단</p><img src="{uri(16, 16)}"><img src="{uri(400, 240)}">'
        f"<p>셋째 문단</p></body></html>"
    ).encode()


def test_html_img_occurrence_mapped():
    items = _enumerate_images_for_classify("페이지.html", _html_with_images())
    # 2번째(16px)는 필터 탈락 — 등장 순번(page)은 1, 3 으로 유지된다.
    assert [(it["index"], it["page"]) for it in items] == [(0, 1), (1, 3)]
    assert all(it["unit"] == "img" for it in items)


def test_html_parse_inserts_markers_at_img():
    from app.ingestion.parsers import parse_document

    text = parse_document("페이지.html", _html_with_images(), "text", insert_image_markers=True)
    lines = text.split("\n")
    i1, i2, i3 = lines.index("첫 문단"), lines.index("둘째 문단"), lines.index("셋째 문단")
    assert any("[[IMG:p1:i0]]" in ln for ln in lines[i1 + 1:i2])   # 1번째 img 자리
    assert any("[[IMG:p3:i1]]" in ln for ln in lines[i2 + 1:i3])   # 3번째 img 자리
    assert "[[IMG:p2" not in text                                   # 필터 탈락분은 마커 없음
    plain = parse_document("페이지.html", _html_with_images(), "text")
    assert "[[IMG" not in plain


# ── HWPX 그림 순번 매핑(Step 3) — 리포 실파일 검증 ───────────────────────────

import pathlib

import pytest

_HWPX_SAMPLE = (
    pathlib.Path(__file__).resolve().parents[2] / "samples" / "hwp" / "3-09월_교육_통합_2023.hwpx"
)


@pytest.mark.skipif(not _HWPX_SAMPLE.exists(), reason="samples/hwp 없음")
def test_hwpx_pic_ordinals_mapped():
    items = _enumerate_images_for_classify("교육.hwpx", _HWPX_SAMPLE.read_bytes())
    assert items, "샘플에 BinData 이미지가 있어야 한다"
    mapped = [it for it in items if isinstance(it["page"], int)]
    assert mapped, "그림 등장 순번이 하나 이상 매핑돼야 한다"
    assert all(it["unit"] == "img" for it in mapped)
    # 순번은 양의 정수이고 (등장 축이므로) 서로 다르다 — 반복 미디어는 첫 순번만.
    pages = [it["page"] for it in mapped]
    assert all(p >= 1 for p in pages) and len(set(pages)) == len(pages)


def _rhwp_with_images_available() -> bool:
    try:
        import rhwp_py

        return hasattr(rhwp_py, "extract_markdown_pages_with_images")
    except ImportError:
        return False


@pytest.mark.skipif(
    not _HWPX_SAMPLE.exists() or not _rhwp_with_images_available(),
    reason="samples/hwp 또는 rhwp_py(이미지 토큰 API) 없음",
)
def test_hwpx_rhwp_parse_inserts_markers():
    """rhwp 전략에서도 이미지 토큰 → 위치 마커 치환(Step 3 — rhwp 앵커 산출)."""
    from app.ingestion.image_classifier import MARKER_RE, media_ref_map
    from app.ingestion.parsers import parse_document

    data = _HWPX_SAMPLE.read_bytes()
    text = parse_document("교육.hwpx", data, "rhwp", insert_image_markers=True)
    assert "[[RHWP_IMAGE:" not in text  # rhwp 토큰은 전부 소진(치환/제거)
    found = {(int(p), int(i)) for p, i in MARKER_RE.findall(text)}
    valid = {(r["page"], r["index"]) for r in media_ref_map("교육.hwpx", data).values()}
    assert found and found <= valid          # 마커는 분류 items 키의 부분집합
    assert len(found) == len(MARKER_RE.findall(text))  # 같은 이미지 반복 치환 없음(첫 등장만)
    # 마커 없이 파싱하면 기존과 동일(토큰 제거 경로).
    plain = parse_document("교육.hwpx", data, "rhwp")
    assert "[[IMG" not in plain and "[[RHWP_IMAGE" not in plain


@pytest.mark.skipif(not _HWPX_SAMPLE.exists(), reason="samples/hwp 없음")
def test_hwpx_parse_inserts_markers():
    from app.ingestion.image_classifier import MARKER_RE, media_image_layout
    from app.ingestion.parsers import parse_document

    data = _HWPX_SAMPLE.read_bytes()
    text = parse_document("교육.hwpx", data, "text", insert_image_markers=True)
    found = {(int(p), int(i)) for p, i in MARKER_RE.findall(text)}
    expected = {
        (it["page"], it["index"])
        for its in media_image_layout("교육.hwpx", data).values()
        for it in its
    }
    assert expected and found == expected  # cap 선별분 전부, 그 외 없음
    plain = parse_document("교육.hwpx", data, "text")
    assert "[[IMG" not in plain


# ── 상한 선별: 큰 이미지 우선 + 원래 순서 유지 ────────────────────────────────

def test_cap_largest_keeps_biggest_in_original_order():
    imgs = [
        {"index": 0, "width": 100, "height": 100},  # 10,000
        {"index": 1, "width": 500, "height": 400},  # 200,000
        {"index": 2, "width": 50, "height": 50},    # 2,500
        {"index": 3, "width": 300, "height": 300},  # 90,000
    ]
    kept = _cap_largest(imgs, 2)
    assert [im["index"] for im in kept] == [1, 3]   # 큰 2장, 원래 순서

    assert _cap_largest(imgs, 0) == imgs            # cap 미설정(0) → 전체
    assert _cap_largest(imgs, 10) == imgs           # 상한 미만 → 전체
